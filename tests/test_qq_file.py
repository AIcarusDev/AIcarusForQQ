from __future__ import annotations

import asyncio
import base64
import json
import os
import sqlite3
from functools import wraps
from pathlib import Path, PurePosixPath
from types import SimpleNamespace

import pytest
from websockets.protocol import State as WsState

from platforms.qq.adapter.client import QQAdapterClient
from platforms.qq.files.logical import collision_name, sanitize_filename
from platforms.qq.files.parsers import ParseContractError, parse_document
from platforms.qq.files.repository import QQFileRepository
from platforms.qq.files.service import QQFileError, QQFileService
from platforms.qq.files.storage import (
    HostFallbackStorage,
    LinuxWorkspaceStorage,
    StorageError,
    StorageRouter,
)
from skills.registry import build_skill_block_for_namespaces
from tools import build_tools
from tools.namespaces import NamespaceRuntimeState, load_namespace_registry
from workspace import WorkspaceService, WslWorkspaceBackend
from workspace.control import detect_workspace_presence


def _async_test(function):
    @wraps(function)
    def run(*args, **kwargs):
        return asyncio.run(function(*args, **kwargs))

    return run


def _create_schema(path: Path) -> None:
    with sqlite3.connect(path) as db:
        db.executescript(
            """
            CREATE TABLE qq_file_downloads (
                download_id TEXT PRIMARY KEY, agent_qq TEXT NOT NULL, session_key TEXT NOT NULL,
                message_id TEXT NOT NULL, conversation_type TEXT NOT NULL, conversation_id TEXT NOT NULL,
                original_filename TEXT NOT NULL, source_file_id TEXT NOT NULL, status TEXT NOT NULL,
                bytes_downloaded INTEGER NOT NULL DEFAULT 0, total_bytes INTEGER, target_path TEXT NOT NULL,
                local_path TEXT, storage_backend TEXT NOT NULL, storage_relpath TEXT NOT NULL,
                failure_code TEXT, failure_message TEXT, failure_retryable INTEGER,
                created_at TEXT NOT NULL, updated_at TEXT NOT NULL, finished_at TEXT
            );
            CREATE UNIQUE INDEX idx_qq_file_downloads_active_source
              ON qq_file_downloads(agent_qq, session_key, message_id)
              WHERE status IN ('queued','resolving','downloading','verifying');
            CREATE TABLE qq_file_records (
                record_id TEXT PRIMARY KEY, agent_qq TEXT NOT NULL, session_key TEXT NOT NULL,
                message_id TEXT NOT NULL, conversation_type TEXT NOT NULL, conversation_id TEXT NOT NULL,
                original_filename TEXT NOT NULL, local_path TEXT NOT NULL, storage_backend TEXT NOT NULL,
                storage_relpath TEXT NOT NULL, size_bytes INTEGER NOT NULL,
                origin TEXT NOT NULL DEFAULT 'qq_download', file_id TEXT NOT NULL DEFAULT '',
                downloaded_at TEXT NOT NULL,
                deleted_at TEXT
            );
            CREATE TABLE qq_file_messages (
                agent_qq TEXT NOT NULL, session_key TEXT NOT NULL, message_id TEXT NOT NULL,
                conversation_type TEXT NOT NULL, conversation_id TEXT NOT NULL, filename TEXT NOT NULL,
                extension TEXT, size_bytes INTEGER, sender_id TEXT NOT NULL DEFAULT '',
                sender_name TEXT NOT NULL DEFAULT '', sent_at TEXT NOT NULL DEFAULT '', indexed_at INTEGER NOT NULL DEFAULT 0,
                PRIMARY KEY(agent_qq, session_key, message_id)
            );
            CREATE TABLE chat_sessions (
                session_key TEXT PRIMARY KEY, focus_name TEXT NOT NULL DEFAULT '', conv_name TEXT NOT NULL DEFAULT ''
            );
            """
        )


class _Router:
    def __init__(self, storage: HostFallbackStorage) -> None:
        self.storage = storage

    def active(self):
        return self.storage

    def frozen(self, backend: str):
        assert backend == self.storage.backend_name
        return self.storage


class _QQClient:
    connected = True
    bot_id = "4242"
    last_bot_id = "4242"

    def __init__(self, payload: bytes = b"hello\nworld\n") -> None:
        self.payload = payload
        self.download_calls = 0

    async def send_api(self, action, params, timeout=15.0):
        assert action == "get_msg"
        return {
            "message_type": "group",
            "group_id": "7777",
            "message": [
                {
                    "type": "file",
                    "data": {
                        "file_id": "encoded-file-id",
                        "file_name": "report.txt",
                        "file_size": len(self.payload),
                    },
                }
            ],
        }

    async def download_file_stream(self, file_id, destination, *, on_progress=None, **_kwargs):
        self.download_calls += 1
        midpoint = len(self.payload) // 2
        with Path(destination).open("wb") as handle:
            for chunk in (self.payload[:midpoint], self.payload[midpoint:]):
                handle.write(chunk)
                if on_progress:
                    result = on_progress(handle.tell(), len(self.payload))
                    if asyncio.iscoroutine(result):
                        await result
        return {"file_name": "report.txt", "size_bytes": len(self.payload)}


class _BlockingQQClient(_QQClient):
    def __init__(self) -> None:
        super().__init__(b"pending")
        self.started = asyncio.Event()

    async def download_file_stream(self, file_id, destination, *, on_progress=None, **_kwargs):
        self.started.set()
        Path(destination).write_bytes(b"part")
        if on_progress:
            result = on_progress(4, len(self.payload))
            if asyncio.iscoroutine(result):
                await result
        await asyncio.Event().wait()


def _service(tmp_path: Path, client=None) -> tuple[QQFileService, HostFallbackStorage]:
    db_path = tmp_path / "qq-file.db"
    _create_schema(db_path)
    storage = HostFallbackStorage(tmp_path)
    service = QQFileService(
        client or _QQClient(),
        None,
        project_root=tmp_path,
        repository=QQFileRepository(str(db_path)),
        storage_router=_Router(storage),
    )
    return service, storage


@_async_test
async def test_restart_recovery_is_scoped_to_the_current_qq_account(tmp_path: Path) -> None:
    database_path = tmp_path / "qq-file.db"
    _create_schema(database_path)
    repository = QQFileRepository(str(database_path))

    async def create(account: str, message: str) -> dict:
        return await repository.create_job({
            "agent_qq": account,
            "session_key": f"qq:private:{account}",
            "message_id": message,
            "conversation_type": "private",
            "conversation_id": account,
            "original_filename": "pending.bin",
            "source_file_id": "source",
            "total_bytes": 1,
            "target_path": f"/home/agent/qq/{account}/file/private_{account}/pending.bin",
            "storage_backend": "host_fallback",
            "storage_relpath": f"qq/{account}/file/private_{account}/pending.bin",
        })

    current = await create("4242", "message-a")
    other = await create("5252", "message-b")
    await repository.recover_interrupted("4242")

    assert (await repository.get_job_row(current["download_id"]))["failure_code"] == "download_interrupted"
    assert (await repository.get_job_row(other["download_id"]))["status"] == "queued"


def test_namespace_is_qq_only_folded_and_loads_skill_on_open() -> None:
    registry = load_namespace_registry()
    client = SimpleNamespace(connected=True, _loop=None, bot_id="4242")
    session = SimpleNamespace(
        key="qq:group:7777",
        conv_type="group",
        conv_id="7777",
        inbound_received_seq=0,
    )
    closed = build_tools(
        {"platforms": {"qq": {"enabled": True}}},
        qq_client=client, workspace_service=object(), session=session, current_platform="qq"
    )
    assert all(key in closed.latent_specs for key in (
        "qq_file.download", "qq_file.read", "qq_file.list_files", "qq_file.search", "qq_file.delete"
    ))
    assert "qq-file" not in build_skill_block_for_namespaces(
        closed.active_namespace_order, closed.namespace_registry
    )

    state = NamespaceRuntimeState()
    state.open("qq_file", registry, 1)
    opened = build_tools(
        {"platforms": {"qq": {"enabled": True}}},
        namespace_state=state, current_round=1, qq_client=client,
        workspace_service=object(), session=session, current_platform="qq",
    )
    assert all(key in opened.active_specs for key in (
        "qq_file.download", "qq_file.read", "qq_file.list_files", "qq_file.search", "qq_file.delete"
    ))
    assert opened.active_specs["qq_file.read"].result_cdata is True
    skill_block = build_skill_block_for_namespaces(opened.active_namespace_order, opened.namespace_registry)
    assert '<skill name="qq-file" from="namespace.qq_file">' in skill_block
    assert "/home/agent/qq/{qq_id}/file/{conversation_type}_{conversation_id}/" in skill_block

    core = build_tools(
        {"platforms": {"qq": {"enabled": True}}},
        namespace_state=state, current_round=1, qq_client=client,
        workspace_service=object(), session=session, current_platform="core",
    )
    assert not any(key.startswith("qq_file.") for key in core.all_specs)

def test_database_initialization_indexes_synchronized_file_messages(tmp_path: Path, monkeypatch) -> None:
    import database

    database_path = tmp_path / "AICQ.db"
    monkeypatch.setattr(database, "DB_PATH", str(database_path))

    async def run() -> None:
        await database.init_db()
        await database.upsert_chat_session("qq:group:7777", "group", "7777", "测试群")
        await database.save_chat_message(
            "qq:group:7777",
            {
                "role": "user",
                "agent_qq": "4242",
                "message_id": "9988",
                "sender_id": "123",
                "sender_name": "发送者",
                "timestamp": "2026-08-31T10:00:00+08:00",
                "content_segments": [
                    {"type": "file", "filename": "资料.PDF", "size_bytes": 42, "file_id": "opaque"}
                ],
            },
        )
        rows = await QQFileRepository(str(database_path)).history_rows("4242")
        assert len(rows) == 1
        assert rows[0]["message_id"] == "9988"
        assert rows[0]["filename"] == "资料.PDF"
        assert rows[0]["extension"] == "pdf"
        assert rows[0]["conversation_name"] == "测试群"

        await database.update_chat_message_id("qq:group:7777", "9988", "9989")
        rows = await QQFileRepository(str(database_path)).history_rows("4242")
        assert [row["message_id"] for row in rows] == ["9989"]

        assert await database.update_chat_message_recalled(
            "9989", "消息已撤回", "2026-08-31T10:01:00+08:00",
            content_segments=[{"type": "note", "text": "消息已撤回"}],
            session_key="qq:group:7777",
        )
        assert await QQFileRepository(str(database_path)).history_rows("4242") == []

    asyncio.run(run())


def test_database_migrates_generated_file_record_columns(tmp_path: Path, monkeypatch) -> None:
    import database

    database_path = tmp_path / "AICQ-old.db"
    with sqlite3.connect(database_path) as db:
        db.execute(
            """CREATE TABLE qq_file_records (
                record_id TEXT PRIMARY KEY, agent_qq TEXT NOT NULL, session_key TEXT NOT NULL,
                message_id TEXT NOT NULL, conversation_type TEXT NOT NULL, conversation_id TEXT NOT NULL,
                original_filename TEXT NOT NULL, local_path TEXT NOT NULL, storage_backend TEXT NOT NULL,
                storage_relpath TEXT NOT NULL, size_bytes INTEGER NOT NULL, downloaded_at TEXT NOT NULL,
                deleted_at TEXT
            )"""
        )
    monkeypatch.setattr(database, "DB_PATH", str(database_path))

    asyncio.run(database.init_db())

    with sqlite3.connect(database_path) as db:
        columns = {row[1]: row for row in db.execute("PRAGMA table_info(qq_file_records)")}
    assert columns["origin"][4] == "'qq_download'"
    assert columns["file_id"][4] == "''"


@_async_test
async def test_download_dedupes_only_while_the_recorded_path_exists(tmp_path: Path) -> None:
    service, storage = _service(tmp_path)
    session = SimpleNamespace(key="qq:group:7777", conv_type="group", conv_id="7777")

    first = await service.start("1001", session)
    assert first["job"]["status"] == "completed"
    logical = PurePosixPath(first["job"]["local_path"])
    assert storage.host_path(logical).read_bytes() == b"hello\nworld\n"

    duplicate = await service.start("1001", session)
    assert duplicate["outcome"] == "already_exists"
    assert service.qq_client.download_calls == 1

    renamed = logical.with_name("renamed.txt")
    storage.host_path(logical).rename(storage.host_path(renamed))
    downloaded_again = await service.start("1001", session)
    assert downloaded_again["job"]["status"] == "completed"
    assert service.qq_client.download_calls == 2
    assert downloaded_again["job"]["local_path"] == str(logical)

    listed = await service.list_files(scope={"type": "all"}, limit=50, cursor=None, session=None)
    assert {item["name"] for item in listed["files"]} == {"renamed.txt", "report.txt"}
    assert {item["managed"] for item in listed["files"]} == {False, True}

    searched = await service.search(
        source="local", query="report", file_types=["txt"], scope={"type": "all"},
        limit=50, cursor=None, session=None,
    )
    assert searched["files"][0]["path"] == str(logical)
    assert searched["files"][0]["match_type"] == "prefix"

    deleted = await service.delete(str(logical))
    assert deleted["deleted"] is True
    assert deleted["was_managed"] is True
    assert not storage.host_path(logical).exists()


@_async_test
async def test_generated_text_is_stored_as_a_session_file_and_linked_after_send(tmp_path: Path) -> None:
    service, storage = _service(tmp_path)
    session = SimpleNamespace(key="qq:group:7777", conv_type="group", conv_id="7777")
    payload = "# 现场文档\n\n正文\n".encode("utf-8")

    stored = await service.store_generated_text(
        content=payload,
        filename="现场文档.md",
        session=session,
    )

    assert stored["origin"] == "agent_generated"
    assert stored["local_path"] == "/home/agent/qq/4242/file/group_7777/现场文档.md"
    assert storage.host_path(PurePosixPath(stored["local_path"])).read_bytes() == payload

    linked = await service.attach_generated_delivery(
        stored["record_id"],
        message_id="556677",
        file_id="remote-file-id",
    )
    assert linked is not None
    assert linked["origin"] == "agent_generated"
    assert linked["message_id"] == "556677"
    assert linked["file_id"] == "remote-file-id"

    read_back = await service.read(
        source={"message_id": "556677"},
        selection=None,
        cursor=None,
        session=session,
    )
    assert read_back["outcome"] == "content"
    assert read_back["content"] == "1\t# 现场文档\n2\t\n3\t正文"

    listed = await service.list_files(
        scope=None,
        limit=50,
        cursor=None,
        session=session,
    )
    assert listed["files"][0]["origin"] == "agent_generated"
    assert listed["files"][0]["source"]["message_id"] == "556677"
    assert listed["files"][0]["source"]["file_id"] == "remote-file-id"


@_async_test
async def test_host_fallback_reversibly_maps_windows_incompatible_names(tmp_path: Path) -> None:
    service, storage = _service(tmp_path)
    logical = PurePosixPath(
        "/home/agent/qq/4242/file/private_123/report:final?.txt"
    )
    source = tmp_path / "source.bin"
    source.write_bytes(b"mapped")
    await storage.commit(source, logical)
    physical = storage.host_path(logical)
    assert physical.name.startswith("~q")
    assert ":" not in physical.name and "?" not in physical.name

    for name in ("~qnotes.txt", "CON.txt"):
        special = logical.with_name(name)
        await storage.commit(source, special)
        assert storage.host_path(special).name.startswith("~q%")
    rows = await storage.list(logical.parent)
    assert {row.path for row in rows} == {
        str(logical), str(logical.with_name("~qnotes.txt")), str(logical.with_name("CON.txt"))
    }


def test_collision_suffix_keeps_linux_filename_within_component_limit() -> None:
    original = sanitize_filename("文" * 200 + ".txt")
    collision = collision_name(original, 1)
    assert len(original.encode("utf-8")) <= 255
    assert len(collision.encode("utf-8")) <= 255
    assert collision.endswith("(1).txt")
    assert sanitize_filename(".aicq-qq-file-user.bin") == "_.aicq-qq-file-user.bin"


def test_storage_router_falls_back_only_when_linux_is_absent(tmp_path: Path, monkeypatch) -> None:
    import platforms.qq.files.storage as storage_module

    router = StorageRouter(tmp_path, object())
    monkeypatch.setattr(storage_module, "detect_workspace_presence", lambda: "absent")
    assert isinstance(router.active(), HostFallbackStorage)

    monkeypatch.setattr(storage_module, "detect_workspace_presence", lambda: "present")
    assert isinstance(router.active(), LinuxWorkspaceStorage)

    monkeypatch.setattr(storage_module, "detect_workspace_presence", lambda: "unknown")
    with pytest.raises(StorageError) as error:
        router.active()
    assert error.value.code == "runtime_unavailable"


@_async_test
async def test_current_scope_requires_a_concrete_qq_session(tmp_path: Path) -> None:
    service, _storage = _service(tmp_path)

    with pytest.raises(QQFileError) as download_error:
        await service.start("1001", None)
    assert getattr(download_error.value, "code", None) == "no_current_qq_session"

    with pytest.raises(QQFileError) as list_error:
        await service.list_files(scope=None, limit=10, cursor=None, session=None)
    assert getattr(list_error.value, "code", None) == "no_current_qq_session"

    with pytest.raises(QQFileError) as search_error:
        await service.search(
            source="history", query="report", file_types=None, scope=None,
            limit=10, cursor=None, session=None,
        )
    assert getattr(search_error.value, "code", None) == "no_current_qq_session"


@_async_test
async def test_timed_out_download_is_discoverable_and_stoppable(tmp_path: Path, monkeypatch) -> None:
    import platforms.qq.files.service as service_module

    client = _BlockingQQClient()
    service, _storage = _service(tmp_path, client)
    monkeypatch.setattr(service_module, "DOWNLOAD_OBSERVATION_SECONDS", 0.01)
    session = SimpleNamespace(key="qq:group:7777", conv_type="group", conv_id="7777")

    started = await service.start("2002", session)
    assert started["observation_timeout"] is True
    assert started["job"]["status"] in {"resolving", "downloading"}
    listed = await service.list_downloads(None, 0, 20)
    assert [job["download_id"] for job in listed["active"]] == [started["job"]["download_id"]]
    stopped = await service.stop(started["job"]["download_id"])
    assert stopped["job"]["status"] == "stopped"
    polled = await service.poll(started["job"]["download_id"])
    assert polled["job"]["status"] == "stopped"


@_async_test
async def test_concurrent_downloads_reserve_distinct_collision_names(tmp_path: Path, monkeypatch) -> None:
    import platforms.qq.files.service as service_module

    service, _storage = _service(tmp_path, _BlockingQQClient())
    monkeypatch.setattr(service_module, "DOWNLOAD_OBSERVATION_SECONDS", 0.01)
    session = SimpleNamespace(key="qq:group:7777", conv_type="group", conv_id="7777")
    first, second = await asyncio.gather(service.start("a", session), service.start("b", session))
    targets = {first["job"]["target_path"], second["job"]["target_path"]}
    assert targets == {
        "/home/agent/qq/4242/file/group_7777/report.txt",
        "/home/agent/qq/4242/file/group_7777/report(1).txt",
    }
    await asyncio.gather(
        service.stop(first["job"]["download_id"]),
        service.stop(second["job"]["download_id"]),
    )


@_async_test
async def test_read_message_returns_download_pending_without_cancelling(tmp_path: Path, monkeypatch) -> None:
    import platforms.qq.files.service as service_module

    client = _BlockingQQClient()
    service, _storage = _service(tmp_path, client)
    monkeypatch.setattr(service_module, "DOWNLOAD_OBSERVATION_SECONDS", 0.01)
    session = SimpleNamespace(key="qq:group:7777", conv_type="group", conv_id="7777")
    result = await service.read(
        source={"message_id": "2112"}, selection=None, cursor=None, session=session
    )
    assert result["outcome"] == "download_pending"
    download_id = result["download"]["download_id"]
    assert (await service.poll(download_id))["job"]["status"] == "downloading"
    await service.stop(download_id)


@_async_test
async def test_history_search_uses_only_the_synchronized_account_index(tmp_path: Path) -> None:
    service, _storage = _service(tmp_path)
    with sqlite3.connect(service.repository.db_path) as db:
        db.execute(
            "INSERT INTO chat_sessions(session_key, focus_name, conv_name) VALUES (?,?,?)",
            ("qq:private:123", "示例好友", "示例好友"),
        )
        db.execute(
            """INSERT INTO qq_file_messages
               (agent_qq,session_key,message_id,conversation_type,conversation_id,filename,extension,
                size_bytes,sender_id,sender_name,sent_at,indexed_at)
               VALUES (?,?,?,?,?,?,?,?,?,?,?,?)""",
            (
                "4242", "qq:private:123", "998877", "private", "123",
                "Ｍonthly-Report.PDF", "pdf", None, "123", "示例好友",
                "2026-08-29T09:00:00+08:00", 1,
            ),
        )
        db.execute(
            """INSERT INTO qq_file_messages
               (agent_qq,session_key,message_id,conversation_type,conversation_id,filename,extension,
                size_bytes,sender_id,sender_name,sent_at,indexed_at)
               VALUES (?,?,?,?,?,?,?,?,?,?,?,?)""",
            (
                "999999", "qq:private:456", "other", "private", "456",
                "Monthly-Report.PDF", "pdf", 1, "456", "另一个账号",
                "2026-08-30T09:00:00+08:00", 1,
            ),
        )
        db.commit()
    result = await service.search(
        source="history", query="monthly-report", file_types=["PDF"], scope={"type": "all"},
        limit=50, cursor=None, session=None,
    )
    assert result["history_coverage"] == "aicq_synced_only"
    assert [message["message_id"] for message in result["messages"]] == ["998877"]
    assert result["messages"][0]["conversation"]["name"] == "示例好友"


@_async_test
async def test_read_text_uses_authenticated_cursor_and_detects_changes(tmp_path: Path) -> None:
    payload = ("一行内容\n" * 1500).encode("utf-8")
    service, storage = _service(tmp_path, _QQClient(payload))
    session = SimpleNamespace(key="qq:group:7777", conv_type="group", conv_id="7777")
    downloaded = await service.start("3003", session)
    path = downloaded["job"]["local_path"]

    first = await service.read(source={"path": path}, selection=None, cursor=None, session=session)
    assert first["outcome"] == "content"
    assert len(first["content"]) == 8000
    assert first["has_more"] is True
    assert first["locations"][0]["ends_mid_unit"] is True

    second = await service.read(source=None, selection=None, cursor=first["next_cursor"], session=session)
    assert second["locations"][0]["starts_mid_unit"] is True

    logical = PurePosixPath(path)
    storage.host_path(logical).write_bytes(payload + b"changed")
    with pytest.raises(Exception) as error:
        await service.read(source=None, selection=None, cursor=first["next_cursor"], session=session)
    assert getattr(error.value, "code", "") == "file_changed"


def test_document_parsers_cover_docx_xlsx_pptx_and_pdf_ocr_boundary(tmp_path: Path) -> None:
    from docx import Document
    from openpyxl import Workbook
    from pptx import Presentation
    from pypdf import PdfWriter

    docx_path = tmp_path / "sample.docx"
    document = Document()
    document.add_heading("标题", level=1)
    document.add_paragraph("正文")
    table = document.add_table(rows=1, cols=2)
    table.cell(0, 0).text = "A"
    table.cell(0, 1).text = "B"
    document.save(docx_path)
    docx_result = parse_document(str(docx_path), None)
    assert docx_result["file_type"] == "docx"
    assert "标题" in docx_result["text"] and "A\tB" in docx_result["text"]

    xlsx_path = tmp_path / "sample.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "汇总"
    sheet["A1"] = "项目"
    sheet["B1"] = 12
    workbook.save(xlsx_path)
    xlsx_result = parse_document(str(xlsx_path), {"type": "xlsx_range", "sheet": "汇总", "cell_range": "A1:B1"})
    assert xlsx_result["file_type"] == "xlsx"
    assert "项目\t12" in xlsx_result["text"]

    pptx_path = tmp_path / "sample.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "演示标题"
    slide.placeholders[1].text = "演示正文"
    presentation.save(pptx_path)
    pptx_result = parse_document(str(pptx_path), None)
    assert pptx_result["file_type"] == "pptx"
    assert "演示标题" in pptx_result["text"]

    pdf_path = tmp_path / "scan.pdf"
    writer = PdfWriter()
    writer.add_blank_page(width=100, height=100)
    with pdf_path.open("wb") as handle:
        writer.write(handle)
    with pytest.raises(ParseContractError) as pdf_error:
        parse_document(str(pdf_path), None)
    assert pdf_error.value.code == "ocr_required"


@_async_test
async def test_adapter_download_stream_treats_file_complete_as_terminal_packet(tmp_path: Path) -> None:
    payload = b"streamed-file-content"
    client = QQAdapterClient()

    class FakeWebSocket:
        state = WsState.OPEN

        async def send(self, raw: str) -> None:
            request = json.loads(raw)
            echo = request["echo"]
            queue = client._api_streams[echo]
            queue.put_nowait({"echo": echo, "stream": "stream-action", "status": "ok", "data": {"type": "stream", "data_type": "file_info", "file_name": "x.bin", "file_size": len(payload)}})
            for index, chunk in enumerate((payload[:7], payload[7:])):
                queue.put_nowait({"echo": echo, "stream": "stream-action", "status": "ok", "data": {"type": "stream", "data_type": "file_chunk", "index": index, "data": base64.b64encode(chunk).decode("ascii"), "size": len(chunk)}})
            queue.put_nowait({"echo": echo, "stream": "stream-action", "status": "ok", "data": {"type": "response", "data_type": "file_complete", "total_bytes": len(payload)}})

    client._ws = FakeWebSocket()
    output = tmp_path / "payload.bin"
    result = await client.download_file_stream("file-id", output, packet_timeout=0.1)
    assert result == {"file_name": "x.bin", "size_bytes": len(payload)}
    assert output.read_bytes() == payload
    assert client._api_streams == {}

    class Sink:
        def __init__(self) -> None:
            self.expected = None
            self.content = bytearray()
            self.finished = False
            self.aborted = False

        async def begin(self, expected_size: int) -> None:
            self.expected = expected_size

        async def write(self, chunk: bytes) -> None:
            self.content.extend(chunk)

        async def finish(self) -> int:
            self.finished = True
            return len(self.content)

        async def abort(self) -> None:
            self.aborted = True

        async def rollback(self) -> None:
            self.aborted = True

    sink = Sink()
    sink_result = await client.download_file_stream("file-id", sink, packet_timeout=0.1)
    assert sink_result["size_bytes"] == len(payload)
    assert bytes(sink.content) == payload
    assert sink.expected == len(payload)
    assert sink.finished is True and sink.aborted is False


@_async_test
async def test_adapter_download_stream_cleans_registration_when_send_fails(tmp_path: Path) -> None:
    client = QQAdapterClient()

    class FailingWebSocket:
        state = WsState.OPEN

        async def send(self, _raw: str) -> None:
            raise ConnectionError("closed")

    client._ws = FailingWebSocket()
    destination = tmp_path / "never-opened.bin"
    with pytest.raises(ConnectionError):
        await client.download_file_stream("file-id", destination)
    assert client._api_streams == {}
    assert not destination.exists()


@pytest.mark.skipif(
    os.environ.get("AICQ_RUN_QQ_FILE_LINUX_INTEGRATION") != "1"
    or detect_workspace_presence() != "present",
    reason="requires an installed AICQ-Workspace and explicit integration opt-in",
)
def test_real_linux_bridge_streams_lists_and_deletes_without_host_staging() -> None:
    async def run() -> None:
        service = WorkspaceService(WslWorkspaceBackend())
        storage = LinuxWorkspaceStorage(service)
        logical = PurePosixPath(f"/home/agent/.aicq/qq-file-integration-{os.getpid()}.bin")
        payload = b"direct-linux-stream"
        try:
            sink = storage.download_sink(logical)
            await sink.begin(len(payload))
            await sink.write(payload[:7])
            await sink.write(payload[7:])
            assert await sink.finish() == len(payload)
            inspected = await storage.stat(logical)
            assert inspected is not None and inspected.kind == "regular"
            assert inspected.size_bytes == len(payload)
            assert any(item.path == str(logical) for item in await storage.list(logical.parent))
            assert await storage.delete(logical) == len(payload)
            assert await storage.stat(logical) is None
        finally:
            try:
                await storage.delete(logical)
            except Exception:
                pass
            await service.close()

    asyncio.run(run())
