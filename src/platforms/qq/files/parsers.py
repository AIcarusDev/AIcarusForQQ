"""Deterministic, non-executing document extraction for QQ files."""

from __future__ import annotations

import re
import zipfile
from pathlib import Path
from typing import Any


MAX_READ_BYTES = 256 * 1024 * 1024
MAX_ARCHIVE_BYTES = 1024 * 1024 * 1024
MAX_ARCHIVE_ENTRIES = 10_000
MAX_EXTRACTED_CHARS = 16 * 1024 * 1024


class ParseContractError(RuntimeError):
    def __init__(self, code: str, message: str, *, details: dict[str, Any] | None = None) -> None:
        super().__init__(message)
        self.code = code
        self.details = details or {}


def _archive_kind(path: Path) -> str | None:
    try:
        with zipfile.ZipFile(path) as archive:
            entries = archive.infolist()
            if len(entries) > MAX_ARCHIVE_ENTRIES:
                raise ParseContractError("archive_safety_limit_exceeded", "文档压缩包条目数量超过安全限制")
            total = 0
            for item in entries:
                name = item.filename.replace("\\", "/")
                if name.startswith("/") or any(part in {".", ".."} for part in name.split("/")):
                    raise ParseContractError("archive_safety_limit_exceeded", "文档压缩包包含不安全路径")
                total += max(0, int(item.file_size))
                if total > MAX_ARCHIVE_BYTES:
                    raise ParseContractError("archive_safety_limit_exceeded", "文档解压后大小超过安全限制")
            content_types = archive.read("[Content_Types].xml").decode("utf-8", errors="replace")
    except zipfile.BadZipFile:
        return None
    except KeyError:
        return None
    if "wordprocessingml.document.main+xml" in content_types:
        return "docx"
    if "spreadsheetml.sheet.main+xml" in content_types:
        return "xlsx"
    if "presentationml.presentation.main+xml" in content_types:
        return "pptx"
    return None


def detect_type(path: Path) -> str:
    extension = path.suffix.casefold().lstrip(".")
    if extension in {"docm", "xlsm", "pptm", "doc", "xls", "ppt"}:
        raise ParseContractError("unsupported_file_type", "该文件类型不能由 qq_file.read 直接读取。", details={"detected_file_type": extension})
    with path.open("rb") as handle:
        sample = handle.read(64 * 1024)
    prefix = sample[:8]
    if prefix.startswith(b"%PDF-"):
        return "pdf"
    if prefix.startswith(b"PK"):
        archive = _archive_kind(path)
        if archive:
            return archive
    text_extensions = {
        "txt", "md", "markdown", "rst", "csv", "tsv", "json", "jsonl", "xml", "yaml", "yml",
        "toml", "ini", "cfg", "conf", "log", "py", "js", "jsx", "ts", "tsx", "css", "html",
        "htm", "sh", "bash", "zsh", "ps1", "bat", "cmd", "c", "h", "cpp", "hpp", "java", "rs",
        "go", "rb", "php", "sql", "tex", "vue", "svelte", "env", "gitignore",
    }
    try:
        decoded_sample = sample.decode("utf-8-sig", errors="strict")
    except UnicodeDecodeError:
        if extension in text_extensions or not extension:
            return "text"
    else:
        binary_controls = sum(
            1 for char in decoded_sample if ord(char) < 32 and char not in {"\n", "\r", "\t", "\f"}
        )
        if "\x00" not in decoded_sample and binary_controls <= max(2, len(decoded_sample) // 100):
            return "text"
    raise ParseContractError(
        "unsupported_file_type",
        "该文件类型不能由 qq_file.read 直接读取。",
        details={"detected_file_type": extension or None},
    )


def _bounds(selection: dict[str, Any] | None, expected: str, total: int, start_key: str, end_key: str) -> tuple[int, int]:
    if selection and selection.get("type") != expected:
        raise ParseContractError("selection_type_mismatch", "选择器类型与文件格式不匹配", details={"selection_type": selection.get("type"), "actual_file_type": expected})
    start = int((selection or {}).get(start_key, 1) or 1)
    end_value = (selection or {}).get(end_key)
    end = int(end_value) if end_value is not None else total
    if start < 1 or end < start:
        raise ParseContractError("invalid_selection", "读取范围无效", details={"selection_type": expected})
    if total == 0:
        return 1, 0
    if start > total:
        raise ParseContractError("invalid_selection", "读取范围超出文档边界", details={"selection_type": expected})
    return start, min(end, total)


def _limit(text: str) -> str:
    if len(text) > MAX_EXTRACTED_CHARS:
        raise ParseContractError("read_limit_exceeded", "所选文档内容超过单次解析复杂度限制")
    return text


def _parse_text(path: Path, selection: dict[str, Any] | None) -> dict[str, Any]:
    try:
        value = path.read_bytes().decode("utf-8-sig", errors="strict")
    except UnicodeDecodeError as exc:
        raise ParseContractError("unsupported_text_encoding", "文本文件不是受支持的 UTF-8 编码") from exc
    lines = value.splitlines()
    start, end = _bounds(selection, "text_lines", len(lines), "start_line", "end_line")
    content = "\n".join(f"{index}\t{lines[index - 1]}" for index in range(start, end + 1))
    return {
        "file_type": "text",
        "document": {"type": "text", "encoding": "utf-8", "total_lines": len(lines)},
        "text": _limit(content),
        "location": {"type": "text_lines", "start_line": start, "end_line": end},
        "warnings": [],
    }


def _parse_pdf(path: Path, selection: dict[str, Any] | None) -> dict[str, Any]:
    from pypdf import PdfReader

    reader = PdfReader(str(path), strict=False)
    if reader.is_encrypted:
        raise ParseContractError("password_required", "PDF 需要密码才能读取")
    start, end = _bounds(selection, "pdf_pages", len(reader.pages), "start_page", "end_page")
    pages: list[str] = []
    empty_pages = 0
    for index in range(start, end + 1):
        text = reader.pages[index - 1].extract_text() or ""
        if not text.strip():
            empty_pages += 1
        pages.append(f"[Page {index}]\n{text.strip()}")
    if pages and empty_pages == len(pages):
        raise ParseContractError("ocr_required", "所选 PDF 页面没有可提取的原生文字，需要 OCR；qq_file.read 暂不支持 OCR。")
    metadata = reader.metadata or {}
    warnings = []
    if empty_pages:
        warnings.append({"code": "partial_ocr_required", "message": "部分 PDF 页面没有可提取的原生文字。"})
    return {
        "file_type": "pdf",
        "document": {
            "type": "pdf",
            "page_count": len(reader.pages),
            "title": str(metadata.get("/Title") or ""),
            "author": str(metadata.get("/Author") or ""),
        },
        "text": _limit("\n\n".join(pages)),
        "location": {"type": "pdf_pages", "start_page": start, "end_page": end},
        "warnings": warnings,
    }


def _parse_docx(path: Path, selection: dict[str, Any] | None) -> dict[str, Any]:
    from docx import Document
    from docx.table import Table
    from docx.text.paragraph import Paragraph
    from docx.oxml.table import CT_Tbl
    from docx.oxml.text.paragraph import CT_P

    document = Document(str(path))
    blocks: list[str] = []
    for child in document.element.body.iterchildren():
        if isinstance(child, CT_P):
            paragraph = Paragraph(child, document)
            style = str(paragraph.style.name or "") if paragraph.style else ""
            label = f" [{style}]" if style else ""
            text = paragraph.text
            links = []
            for hyperlink in getattr(paragraph, "hyperlinks", ()):
                url = str(getattr(hyperlink, "url", "") or "")
                label_text = str(getattr(hyperlink, "text", "") or "")
                if url:
                    links.append(f"[Link: {label_text} -> {url}]")
            if links:
                text = text + ("\n" if text else "") + "\n".join(links)
            blocks.append(f"[Block {len(blocks) + 1}{label}]\n{text}")
        elif isinstance(child, CT_Tbl):
            table = Table(child, document)
            rows = ["\t".join(cell.text for cell in row.cells) for row in table.rows]
            blocks.append(f"[Block {len(blocks) + 1} | Table]\n" + "\n".join(rows))
    start, end = _bounds(selection, "docx_blocks", len(blocks), "start_block", "end_block")
    warnings = []
    if document.inline_shapes:
        warnings.append({"code": "embedded_media_omitted", "message": "DOCX 中的嵌入媒体未作为正文读取。"})
    return {
        "file_type": "docx",
        "document": {"type": "docx", "block_count": len(blocks)},
        "text": _limit("\n\n".join(blocks[start - 1 : end])),
        "location": {"type": "docx_blocks", "start_block": start, "end_block": end},
        "warnings": warnings,
    }


_A1_RANGE_RE = re.compile(r"^[A-Za-z]{1,3}[1-9][0-9]*(?::[A-Za-z]{1,3}[1-9][0-9]*)?$")


def _parse_xlsx(path: Path, selection: dict[str, Any] | None) -> dict[str, Any]:
    from openpyxl import load_workbook
    from openpyxl.utils.cell import get_column_letter, range_boundaries

    if selection and selection.get("type") != "xlsx_range":
        raise ParseContractError("selection_type_mismatch", "选择器类型与文件格式不匹配", details={"selection_type": selection.get("type"), "actual_file_type": "xlsx_range"})
    workbook = load_workbook(str(path), read_only=True, data_only=False, keep_links=False)
    cached_workbook = load_workbook(str(path), read_only=True, data_only=True, keep_links=False)
    try:
        sheets = workbook.worksheets
        if not sheets:
            selected = []
        elif selection:
            sheet_name = str(selection.get("sheet") or "")
            if sheet_name not in workbook.sheetnames:
                raise ParseContractError("invalid_selection", "指定工作表不存在", details={"selection_type": "xlsx_range"})
            selected = [workbook[sheet_name]]
        else:
            selected = list(sheets)
        pieces: list[str] = []
        warnings: list[dict[str, str]] = []
        missing_formula_cache = False
        locations: list[dict[str, Any]] = []
        for sheet_index, sheet in enumerate(sheets, start=1):
            if sheet not in selected:
                continue
            if selection and selection.get("cell_range"):
                cell_range = str(selection["cell_range"]).upper()
                if not _A1_RANGE_RE.fullmatch(cell_range):
                    raise ParseContractError("invalid_selection", "XLSX 单元格区域不是有效 A1 范围")
                min_col, min_row, max_col, max_row = range_boundaries(cell_range)
            else:
                min_col, min_row = 1, 1
                max_col, max_row = max(1, sheet.max_column), max(1, sheet.max_row)
                cell_range = f"A1:{get_column_letter(max_col)}{max_row}"
            rows: list[str] = []
            cached_sheet = cached_workbook[sheet.title]
            header = "\t" + "\t".join(get_column_letter(column) for column in range(min_col, max_col + 1))
            rows.append(header)
            for cells in sheet.iter_rows(min_row=min_row, max_row=max_row, min_col=min_col, max_col=max_col):
                values = []
                for cell in cells:
                    value = cell.value
                    if isinstance(value, str) and value.startswith("="):
                        cached = cached_sheet[cell.coordinate].value
                        if cached is None:
                            missing_formula_cache = True
                            values.append(value)
                        else:
                            values.append(f"{value} [cached: {cached}]")
                    else:
                        values.append("" if value is None else str(value))
                rows.append(f"{cells[0].row}\t" + "\t".join(values))
            pieces.append(f"[Sheet: {sheet.title} | {cell_range}]\n" + "\n".join(rows))
            locations.append(
                {
                    "type": "xlsx_range",
                    "sheet": sheet.title,
                    "sheet_index": sheet_index,
                    "sheet_state": sheet.sheet_state,
                    "cell_range": cell_range,
                }
            )
        if not locations:
            locations = [{"type": "xlsx_range", "sheet": "", "sheet_index": 0, "sheet_state": "visible", "cell_range": "A1:A1"}]
        if missing_formula_cache:
            warnings.append({"code": "formula_value_unavailable", "message": "部分公式没有可用的缓存值，未执行重新计算。"})
        with zipfile.ZipFile(path) as archive:
            if any(name.startswith("xl/externalLinks/") for name in archive.namelist()):
                warnings.append({"code": "external_link_omitted", "message": "XLSX 外部数据连接未读取或刷新。"})
        return {
            "file_type": "xlsx",
            "document": {"type": "xlsx", "sheet_count": len(sheets)},
            "text": _limit("\n\n".join(pieces)),
            "locations": locations,
            "warnings": warnings,
        }
    finally:
        workbook.close()
        cached_workbook.close()


def _parse_pptx(path: Path, selection: dict[str, Any] | None) -> dict[str, Any]:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    presentation = Presentation(str(path))
    start, end = _bounds(selection, "pptx_slides", len(presentation.slides), "start_slide", "end_slide")
    pieces: list[str] = []
    warnings: list[dict[str, str]] = []
    omitted_media = False
    for index in range(start, end + 1):
        slide = presentation.slides[index - 1]
        texts: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                value = shape.text.strip()
                if value:
                    texts.append(value)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    texts.append("\t".join(cell.text for cell in row.cells))
            if shape.shape_type in {MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.MEDIA} or getattr(shape, "has_chart", False):
                omitted_media = True
        try:
            notes = slide.notes_slide.notes_text_frame.text.strip()
        except Exception:
            notes = ""
        if notes:
            texts.append("[Notes]\n" + notes)
        pieces.append(f"[Slide {index}]\n" + "\n".join(texts))
    if omitted_media:
        warnings.append({"code": "embedded_media_omitted", "message": "PPTX 中的图片、媒体或图表视觉内容未读取。"})
    return {
        "file_type": "pptx",
        "document": {"type": "pptx", "slide_count": len(presentation.slides)},
        "text": _limit("\n\n".join(pieces)),
        "location": {"type": "pptx_slides", "start_slide": start, "end_slide": end},
        "warnings": warnings,
    }


def parse_document(path_value: str, selection: dict[str, Any] | None) -> dict[str, Any]:
    path = Path(path_value)
    size = path.stat().st_size
    if size > MAX_READ_BYTES:
        raise ParseContractError(
            "file_too_large_to_read",
            "文件超过直接读取大小限制",
            details={"size_bytes": size, "limit_bytes": MAX_READ_BYTES},
        )
    kind = detect_type(path)
    try:
        if kind == "text":
            return _parse_text(path, selection)
        if kind == "pdf":
            return _parse_pdf(path, selection)
        if kind == "docx":
            return _parse_docx(path, selection)
        if kind == "xlsx":
            return _parse_xlsx(path, selection)
        if kind == "pptx":
            return _parse_pptx(path, selection)
    except ParseContractError:
        raise
    except Exception as exc:
        raise ParseContractError("parse_failed", "文档解析失败") from exc
    raise ParseContractError("unsupported_file_type", "该文件类型不能由 qq_file.read 直接读取。")


def parse_document_safe(path_value: str, selection: dict[str, Any] | None) -> dict[str, Any]:
    """Pickle-safe process-pool boundary."""

    import socket

    original_socket = socket.socket
    original_create_connection = socket.create_connection

    def blocked_network(*_args, **_kwargs):
        raise OSError("network access is disabled while reading QQ files")

    socket.socket = blocked_network  # type: ignore[assignment]
    socket.create_connection = blocked_network  # type: ignore[assignment]
    try:
        return {"ok": True, "result": parse_document(path_value, selection)}
    except ParseContractError as exc:
        return {
            "ok": False,
            "error": {"code": exc.code, "message": str(exc), "details": exc.details},
        }
    except Exception:
        return {
            "ok": False,
            "error": {"code": "parse_failed", "message": "文档解析失败", "details": {}},
        }
    finally:
        socket.socket = original_socket  # type: ignore[assignment]
        socket.create_connection = original_create_connection  # type: ignore[assignment]
